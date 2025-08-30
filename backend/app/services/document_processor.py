"""
Document processing service for extracting content from various file formats
"""

import asyncio
import logging
import mimetypes
from pathlib import Path
from typing import Dict, Any, Optional, List
from datetime import datetime
import uuid

import fitz  # PyMuPDF
from pptx import Presentation
import openpyxl
from docx import Document as DocxDocument
from PIL import Image
import io

from app.core.config import settings
from app.core.logging import get_logger

logger = get_logger(__name__)


class DocumentProcessingError(Exception):
    """Custom exception for document processing errors"""
    pass


class DocumentProcessor:
    """Main document processor for handling multiple file formats"""
    
    def __init__(self):
        self.supported_formats = {
            'application/pdf': self._process_pdf,
            'application/vnd.ms-powerpoint': self._process_ppt,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._process_pptx,
            'application/vnd.ms-excel': self._process_excel,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._process_excel,
            'application/msword': self._process_word,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._process_word,
        }
    
    async def process_document(
        self, 
        file_path: str, 
        file_type: str,
        document_id: str
    ) -> Dict[str, Any]:
        """
        Process a document and extract content based on file type
        
        Args:
            file_path: Path to the document file
            file_type: MIME type of the document
            document_id: Unique identifier for the document
            
        Returns:
            Dictionary containing extracted content and metadata
        """
        logger.info(f"Starting document processing for {document_id}", 
                   document_id=document_id, file_type=file_type)
        
        start_time = datetime.utcnow()
        
        try:
            # Validate file exists
            if not Path(file_path).exists():
                raise DocumentProcessingError(f"File not found: {file_path}")
            
            # Check if file type is supported
            if file_type not in self.supported_formats:
                raise DocumentProcessingError(f"Unsupported file type: {file_type}")
            
            # Get the appropriate processor
            processor = self.supported_formats[file_type]
            
            # Process the document
            result = await processor(file_path, document_id)
            
            # Add processing metadata
            processing_time = (datetime.utcnow() - start_time).total_seconds()
            result.update({
                'processing_metadata': {
                    'document_id': document_id,
                    'file_path': file_path,
                    'file_type': file_type,
                    'processing_time_seconds': processing_time,
                    'processed_at': datetime.utcnow().isoformat(),
                    'processor_version': '1.0.0'
                }
            })
            
            logger.info(f"Document processing completed for {document_id}", 
                       document_id=document_id, processing_time=processing_time)
            
            return result
            
        except Exception as e:
            logger.error(f"Document processing failed for {document_id}: {str(e)}", 
                        document_id=document_id, error=str(e))
            raise DocumentProcessingError(f"Processing failed: {str(e)}")
    
    async def _process_pdf(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process PDF documents using PyMuPDF"""
        logger.info(f"Processing PDF document {document_id}")
        
        try:
            doc = fitz.open(file_path)
            
            # Extract basic metadata
            metadata = doc.metadata
            page_count = len(doc)
            
            # Extract text content
            full_text = ""
            pages_content = []
            tables = []
            images = []
            
            for page_num in range(page_count):
                page = doc[page_num]
                
                # Extract text
                page_text = page.get_text()
                full_text += page_text + "\n"
                
                pages_content.append({
                    'page_number': page_num + 1,
                    'text': page_text,
                    'char_count': len(page_text)
                })
                
                # Extract tables (basic implementation)
                page_tables = self._extract_pdf_tables(page)
                if page_tables:
                    tables.extend(page_tables)
                
                # Extract images
                page_images = self._extract_pdf_images(page, page_num)
                if page_images:
                    images.extend(page_images)
            
            doc.close()
            
            # Calculate quality metrics
            text_quality = self._calculate_text_quality(full_text)
            
            return {
                'extracted_content': {
                    'full_text': full_text.strip(),
                    'pages': pages_content,
                    'tables': tables,
                    'images': images
                },
                'document_metadata': {
                    'title': metadata.get('title', ''),
                    'author': metadata.get('author', ''),
                    'subject': metadata.get('subject', ''),
                    'creator': metadata.get('creator', ''),
                    'producer': metadata.get('producer', ''),
                    'creation_date': metadata.get('creationDate', ''),
                    'modification_date': metadata.get('modDate', ''),
                    'page_count': page_count
                },
                'quality_metrics': {
                    'text_quality_score': text_quality,
                    'total_characters': len(full_text),
                    'total_words': len(full_text.split()),
                    'tables_found': len(tables),
                    'images_found': len(images)
                }
            }
            
        except Exception as e:
            logger.error(f"PDF processing error for {document_id}: {str(e)}")
            raise DocumentProcessingError(f"PDF processing failed: {str(e)}")
    
    async def _process_pptx(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process PowerPoint presentations"""
        logger.info(f"Processing PowerPoint document {document_id}")
        
        try:
            prs = Presentation(file_path)
            
            # Extract slides content
            slides_content = []
            full_text = ""
            images = []
            charts = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = ""
                slide_images = []
                slide_charts = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                    
                    # Check for images
                    if shape.shape_type == 13:  # Picture
                        slide_images.append({
                            'slide_number': slide_num + 1,
                            'shape_id': shape.shape_id,
                            'alt_text': getattr(shape, 'alt_text', ''),
                        })
                    
                    # Check for charts
                    if hasattr(shape, 'chart'):
                        slide_charts.append({
                            'slide_number': slide_num + 1,
                            'chart_type': str(shape.chart.chart_type),
                            'title': getattr(shape.chart, 'chart_title', ''),
                        })
                
                slides_content.append({
                    'slide_number': slide_num + 1,
                    'text': slide_text.strip(),
                    'char_count': len(slide_text),
                    'images_count': len(slide_images),
                    'charts_count': len(slide_charts)
                })
                
                full_text += slide_text
                images.extend(slide_images)
                charts.extend(slide_charts)
            
            # Calculate quality metrics
            text_quality = self._calculate_text_quality(full_text)
            
            return {
                'extracted_content': {
                    'full_text': full_text.strip(),
                    'slides': slides_content,
                    'images': images,
                    'charts': charts
                },
                'document_metadata': {
                    'slide_count': len(prs.slides),
                    'slide_width': prs.slide_width,
                    'slide_height': prs.slide_height
                },
                'quality_metrics': {
                    'text_quality_score': text_quality,
                    'total_characters': len(full_text),
                    'total_words': len(full_text.split()),
                    'images_found': len(images),
                    'charts_found': len(charts)
                }
            }
            
        except Exception as e:
            logger.error(f"PowerPoint processing error for {document_id}: {str(e)}")
            raise DocumentProcessingError(f"PowerPoint processing failed: {str(e)}")
    
    async def _process_ppt(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process legacy PowerPoint files (redirect to pptx processor)"""
        return await self._process_pptx(file_path, document_id)
    
    async def _process_excel(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process Excel spreadsheets"""
        logger.info(f"Processing Excel document {document_id}")
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            # Extract worksheets content
            worksheets_content = []
            financial_data = []
            tables = []
            full_text = ""
            
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                
                # Extract data from worksheet
                sheet_data = []
                sheet_text = ""
                
                for row in worksheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_data = [str(cell) if cell is not None else "" for cell in row]
                        sheet_data.append(row_data)
                        sheet_text += " ".join(row_data) + "\n"
                
                # Identify potential financial data
                financial_indicators = self._identify_financial_data(sheet_data, sheet_name)
                if financial_indicators:
                    financial_data.extend(financial_indicators)
                
                # Identify tables
                if len(sheet_data) > 1:  # Has header and data
                    tables.append({
                        'sheet_name': sheet_name,
                        'rows': len(sheet_data),
                        'columns': len(sheet_data[0]) if sheet_data else 0,
                        'data_preview': sheet_data[:5]  # First 5 rows
                    })
                
                worksheets_content.append({
                    'sheet_name': sheet_name,
                    'text': sheet_text.strip(),
                    'row_count': len(sheet_data),
                    'column_count': len(sheet_data[0]) if sheet_data else 0,
                    'has_financial_data': len(financial_indicators) > 0
                })
                
                full_text += sheet_text
            
            workbook.close()
            
            # Calculate quality metrics
            text_quality = self._calculate_text_quality(full_text)
            
            return {
                'extracted_content': {
                    'full_text': full_text.strip(),
                    'worksheets': worksheets_content,
                    'financial_data': financial_data,
                    'tables': tables
                },
                'document_metadata': {
                    'worksheet_count': len(workbook.sheetnames),
                    'worksheet_names': workbook.sheetnames
                },
                'quality_metrics': {
                    'text_quality_score': text_quality,
                    'total_characters': len(full_text),
                    'financial_indicators_found': len(financial_data),
                    'tables_found': len(tables)
                }
            }
            
        except Exception as e:
            logger.error(f"Excel processing error for {document_id}: {str(e)}")
            raise DocumentProcessingError(f"Excel processing failed: {str(e)}")
    
    async def _process_word(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process Word documents"""
        logger.info(f"Processing Word document {document_id}")
        
        try:
            doc = DocxDocument(file_path)
            
            # Extract paragraphs
            paragraphs = []
            full_text = ""
            
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append({
                        'text': para.text,
                        'style': para.style.name if para.style else 'Normal'
                    })
                    full_text += para.text + "\n"
            
            # Extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                tables.append({
                    'rows': len(table_data),
                    'columns': len(table_data[0]) if table_data else 0,
                    'data_preview': table_data[:3]  # First 3 rows
                })
            
            # Extract document properties
            core_props = doc.core_properties
            
            # Calculate quality metrics
            text_quality = self._calculate_text_quality(full_text)
            
            return {
                'extracted_content': {
                    'full_text': full_text.strip(),
                    'paragraphs': paragraphs,
                    'tables': tables
                },
                'document_metadata': {
                    'title': core_props.title or '',
                    'author': core_props.author or '',
                    'subject': core_props.subject or '',
                    'created': core_props.created.isoformat() if core_props.created else '',
                    'modified': core_props.modified.isoformat() if core_props.modified else '',
                    'paragraph_count': len(paragraphs),
                    'table_count': len(tables)
                },
                'quality_metrics': {
                    'text_quality_score': text_quality,
                    'total_characters': len(full_text),
                    'total_words': len(full_text.split()),
                    'tables_found': len(tables)
                }
            }
            
        except Exception as e:
            logger.error(f"Word processing error for {document_id}: {str(e)}")
            raise DocumentProcessingError(f"Word processing failed: {str(e)}")
    
    def _extract_pdf_tables(self, page) -> List[Dict[str, Any]]:
        """Extract tables from PDF page (basic implementation)"""
        try:
            # This is a simplified table extraction
            # In production, you might want to use libraries like tabula-py or camelot
            tables = page.find_tables()
            extracted_tables = []
            
            for table in tables:
                table_data = table.extract()
                if table_data:
                    extracted_tables.append({
                        'rows': len(table_data),
                        'columns': len(table_data[0]) if table_data else 0,
                        'data_preview': table_data[:3]  # First 3 rows
                    })
            
            return extracted_tables
        except:
            return []
    
    def _extract_pdf_images(self, page, page_num: int) -> List[Dict[str, Any]]:
        """Extract images from PDF page"""
        try:
            image_list = page.get_images()
            images = []
            
            for img_index, img in enumerate(image_list):
                images.append({
                    'page_number': page_num + 1,
                    'image_index': img_index,
                    'width': img[2] if len(img) > 2 else None,
                    'height': img[3] if len(img) > 3 else None,
                })
            
            return images
        except:
            return []
    
    def _identify_financial_data(self, sheet_data: List[List[str]], sheet_name: str) -> List[Dict[str, Any]]:
        """Identify potential financial data in Excel sheets"""
        financial_indicators = []
        
        # Common financial keywords
        financial_keywords = [
            'revenue', 'income', 'profit', 'loss', 'cash', 'flow',
            'balance', 'assets', 'liabilities', 'equity', 'expenses',
            'cost', 'margin', 'ebitda', 'roi', 'irr', 'npv'
        ]
        
        for row_idx, row in enumerate(sheet_data):
            for col_idx, cell in enumerate(row):
                if isinstance(cell, str):
                    cell_lower = cell.lower()
                    for keyword in financial_keywords:
                        if keyword in cell_lower:
                            financial_indicators.append({
                                'sheet_name': sheet_name,
                                'row': row_idx + 1,
                                'column': col_idx + 1,
                                'keyword': keyword,
                                'context': cell,
                                'value': row[col_idx + 1] if col_idx + 1 < len(row) else None
                            })
        
        return financial_indicators
    
    def _calculate_text_quality(self, text: str) -> float:
        """Calculate text quality score based on various factors"""
        if not text or len(text.strip()) == 0:
            return 0.0
        
        # Basic quality metrics
        char_count = len(text)
        word_count = len(text.split())
        
        # Check for readable content (not just symbols/numbers)
        readable_chars = sum(1 for c in text if c.isalpha() or c.isspace())
        readability_ratio = readable_chars / char_count if char_count > 0 else 0
        
        # Check average word length (reasonable words)
        avg_word_length = char_count / word_count if word_count > 0 else 0
        word_length_score = min(avg_word_length / 6, 1.0)  # Normalize to 0-1
        
        # Combine metrics (simple weighted average)
        quality_score = (readability_ratio * 0.6 + word_length_score * 0.4)
        
        return min(quality_score, 1.0)


# Global document processor instance
document_processor = DocumentProcessor()
